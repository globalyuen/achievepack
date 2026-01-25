from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_presentation():
    prs = Presentation()

    # Define brand colors
    PRIMARY_COLOR = RGBColor(26, 77, 46) 
    SECONDARY_COLOR = RGBColor(240, 253, 244)
    TEXT_COLOR = RGBColor(40, 40, 40)
    
    # Image Paths (using the generated artifact paths)
    HERO_IMAGE = "/Users/ryanmacmini/.gemini/antigravity/brain/1ff3d642-be1a-4c8e-8f9e-ca4976fa5445/ppt_hero_packaging_1769248866451.png"
    PRODUCT_IMAGE = "/Users/ryanmacmini/.gemini/antigravity/brain/1ff3d642-be1a-4c8e-8f9e-ca4976fa5445/ppt_product_showcase_1769248882521.png"

    def add_slide(title_text, content_items, image_path=None):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
        title.text_frame.paragraphs[0].font.bold = True

        # Content
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        
        # If image is present, resize text box to left half
        if image_path and os.path.exists(image_path):
            body_shape.width = Inches(5)
            
            # Add image on the right
            left = Inches(5.5)
            top = Inches(2)
            height = Inches(4.5)
            slide.shapes.add_picture(image_path, left, top, height=height)

        for item in content_items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(20)
            p.level = 0
            p.space_after = Pt(10)
            
    # Slide 1: Title with Hero Image background or side
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Achieve Pack"
    subtitle.text = "Certified Eco-Friendly Packaging Solutions\nfor Global Brands"
    
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Pt(44)
    
    # Add hero image to title slide
    if os.path.exists(HERO_IMAGE):
        left = Inches(1)
        top = Inches(4.5)
        width = Inches(8)
        slide.shapes.add_picture(HERO_IMAGE, left, top, width=width)

    # Slide 2: Our Mission
    add_slide("Our Mission", [
        "Making sustainable packaging the easy choice for every business.",
        "Founded in 2011, serving 500+ brands worldwide.",
        "We believe environmental responsibility shouldn't be a premium feature—it should be the standard.",
        "Bridging the gap between eco-friendly materials and high-performance needs."
    ])
    
    # Slide 3: The Advantage
    add_slide("Why Brands Choose Us", [
        "Low MOQs: Start with just 100 pieces.",
        "Fast Turnaround: Samples in 2-3 weeks.",
        "Real Certifications: EN13432, GRS-Certified.",
        "Cost Effective: Save 15-22% on shipping.",
        "No Greenwashing: Full transparency."
    ], image_path=PRODUCT_IMAGE) # Using product image here

    # Slide 4: Our Solutions
    add_slide("Sustainable Material Options", [
        "Compostable (Kraft + PLA / NatureFlex):",
        "   - Breaks down in 90-180 days.",
        "   - Ideal for coffee, tea, and dry goods.",
        "Recyclable Mono-PE & Mono-PP:",
        "   - Single material structure = Actually recyclable.",
        "   - High barrier options available (EVOH).",
        "PCR (Post-Consumer Recycled):",
        "   - Made from plastics that already existed.",
        "   - GRS Certified content (30-100%)."
    ])

    # Slide 5: Product Formats
    add_slide("Packaging Formats", [
        "Stand-Up Pouches: Great shelf presence.",
        "Flat Bottom Bags: Premium 'box' look.",
        "Spout Pouches: Liquid-ready.",
        "Vacuum Pouches: Extended shelf life.",
        "Rollstock: For automated lines."
    ], image_path=HERO_IMAGE)

    # ... (Rest of existing slides code) ...
    # Slide 6: Industries We Serve
    add_slide("Industries & Applications", [
        "Coffee & Tea: High barrier valves, kraft aesthetics.",
        "Pet Food: Odor-blocking, durable, large formats.",
        "Supplements: Moisture control for powders.",
        "Beauty & Cosmetics: Premium finishes.",
        "Food & Snacks: FDA-compliant, grease-resistant.",
        "Frozen Food: Cold-resistant materials."
    ])

    # Slide 7: Case Studies
    add_slide("Success Stories", [
        "Bean & Bole Coffee: Compostable bags + freshman protection.",
        "VitalGreen Superfoods: Recyclable Mono-PE for 8 SKUs.",
        "Pawsome Naturals: Child-Resistant PCR packaging.",
        "Artisan Cocoa (Dubai): Heat-resistant packaging."
    ], image_path=PRODUCT_IMAGE)
    
    # Slide 8: Technical Excellence
    add_slide("Quality & Certifications", [
        "Certified Compostable: EN 13432, ASTM D6400.",
        "Start-of-Life: FSC Certified Paper, Bio-based PE.",
        "End-of-Life: TUV Austria OK Compost, GRS.",
        "Food Safety: FDA 21 CFR, EU 10/2011 Compliance.",
        "Testing: Seal strength, OTR/MVTR (Barrier checking)."
    ])

    # Slide 9: Contact
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Ready to Switch?"
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    body = slide.placeholders[1]
    tf = body.text_frame
    p = tf.add_paragraph()
    p.text = "Start your sustainable packaging journey today."
    p.font.size = Pt(24)
    p.font.bold = True
    p.space_after = Pt(20)

    p = tf.add_paragraph()
    p.text = "Web: achievepack.com\nEmail: ryan@achievepack.com"
    p.font.size = Pt(20)

    prs.save('AchievePack_Presentation_With_Images.pptx')
    print("Presentation saved as AchievePack_Presentation_With_Images.pptx")

if __name__ == "__main__":
    create_presentation()
